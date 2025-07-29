             
import pyttsx3
from gesture_control import start_gesture_control
import speech_recognition
import requests
from bs4 import BeautifulSoup
import datetime
import pyautogui
import os
import my_keyboard
import random
import webbrowser
import NewsRead
import pywhatkit as kit
import time
import subprocess
import pyautogui
from INTRO import play_gif_with_menu
from dotenv import load_dotenv
from fpdf import FPDF
from scroll_control import scroll_up, scroll_down
from battery import check_battery
from generate_img import generate_image
from handle_system import handle_system_command
from youtube_automation import control_youtube
from handle_system import handle_system_command
from whatsapp_automation import send_whatsapp_message_auto
from pptx import Presentation
from pptx.util import Inches
from docx import Document
import os
import openai
import email_assistant
import calendar_assistant
import reminders_assistant
from advanced_features import get_advanced_features
from advanced_notifications import get_notification_system
from ai_personality import get_personality_manager
from analytics_dashboard import get_analytics_dashboard
from voice_games import get_voice_games
from ai_enhanced_features import get_enhanced_ai
load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")
play_gif_with_menu()

# Initialize advanced features
advanced_features = get_advanced_features()
notification_system = get_notification_system()
personality_manager = get_personality_manager()
analytics_dashboard = get_analytics_dashboard()
voice_games = get_voice_games()
enhanced_ai = get_enhanced_ai()

# Link advanced features with notification system
advanced_features.set_notification_system(notification_system)

engine = pyttsx3.init("sapi5")
voices = engine.getProperty("voices")
engine.setProperty("voice", voices[0].id)
engine.setProperty("rate", 300) 



def speak(audio):
    engine.say(audio)
    engine.runAndWait()


def check_password():
    """Function to verify password before starting Jarvis"""
    try:
        with open("password.txt", "r") as pw_file:
            pw = pw_file.read().strip()  # Remove spaces/newlines
            print

        for i in range(3):  # Allow 3 attempts
            user_pw = input("Enter Password to open Jarvis: ").strip()

            if user_pw == pw:
                print("WELCOME SIR! PLZ SPEAK [WAKE UP] TO LOAD ME UP")
                return True
            elif i == 2:
                print("Too many incorrect attempts. Exiting...")
                exit()
            else:
                print("Incorrect Password. Try Again.")

    except FileNotFoundError:
        print("Error: password.txt not found. Please create the file.")
        exit()




def takeCommand():
    r = speech_recognition.Recognizer()
    with speech_recognition.Microphone() as source:
        print("Listening.....")
        r.pause_threshold = 0.5

        r.energy_threshold = 100
        audio = r.listen(source, 0, 4)

    try:
        print("Understanding..")
        query = r.recognize_google(audio, language='en-in')
        print(f"You Said: {query}\n")
        return query.lower()  # Ensure lowercase for consistency
    except Exception as e:
        print("Say that again")
        return "None"
def chat_with_gpt(prompt):
    system_role = "You are Jarvis, a smart, professional, and friendly desktop voice assistant created by Saheli. You can assist with any task, answer questions, or provide helpful information."

    try:
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": system_role},
                {"role": "user", "content": prompt}
            ],
            max_tokens=200,
            temperature=0.7
        )
        reply = response['choices'][0]['message']['content']
        return reply.strip()
    except Exception as e:
        return "Sorry, I couldn't connect to ChatGPT at the moment."
def create_pdf(text, filename="chatbot_output.pdf"):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.multi_cell(200, 10, txt=text)
    pdf.output(filename)
    speak(f"PDF saved as {filename}")

def create_word_doc(text, filename="chatbot_output.docx"):
    doc = Document()
    doc.add_paragraph(text)
    doc.save(filename)
    speak(f"Word document saved as {filename}")

def create_ppt(text, filename="chatbot_output.pptx"):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Jarvis Presentation"
    content.text = text
    prs.save(filename)
    speak(f"Presentation saved as {filename}")
# alarm set 
def alarm(query):
     timehere = open("Alarmtext.txt","a")
     timehere.write(query)
     timehere.close()
     os.startfile("alarm.py")
# weather report
def get_weather(city="Delhi"):  # Default city is Delhi
    api_key = "cf9b51d9e6eaa91fd1a10945e1c928d9"  # Replace with your OpenWeatherMap API key
    base_url = "http://api.openweathermap.org/data/2.5/weather?"
    complete_url = f"{base_url}q={city}&appid={api_key}&units=metric"

    response = requests.get(complete_url)
    data = response.json()

    if data["cod"] == 200:  # ✅ Check if the request was successful
        temp = data["main"]["temp"]
        description = data["weather"][0]["description"]
        speak(f"The current temperature in {city} is {temp} degrees Celsius with {description}.")
    else:
        speak("Sorry, I couldn't fetch the weather details.")

if __name__ == "__main__": 
     if check_password(): 
        # Start notification service
        reminders_assistant.start_notification_service()
     while True:  
        query = takeCommand()
        # if face_login():  # Only activate Jarvis if face is recognized
        #  speak("Jarvis is ready to assist you.")
        if "wake up" in query:
            from GreetMe import greetMe
            greetMe()

            while True:
                query = takeCommand()

                if "go to sleep" in query:  # ✅ Fixed condition
                    speak("Ok sir, You can call me anytime")
                    print("Sleeping mode activated.")  
                    break  # Exits only after speaking    

                elif "click my photo" in query:
                        pyautogui.press("win")
                        time.sleep(1)
                        pyautogui.write("camera", interval=0.1)
                        time.sleep(1)
                        pyautogui.press("enter")
                        time.sleep(3)
                        speak("SMILE")
                        time.sleep(1)
                        pyautogui.press("enter")
                
                elif "hello" in query:
                    speak("Hello sir, how are you?")
                elif "i am fine" in query:
                    speak("That's great, sir")
                elif "how are you" in query:
                    speak("Perfect, sir")
                elif "thank you" in query:
                    speak("You're welcome, sir")
                elif "tired" in query:
                   speak("Playing your favourite songs, sir")
                   a = (1,2,3)
                   b = random.choice(a)
                   if b==1:
                    webbrowser.open("https://www.youtube.com/watch?v=TVbI55pDdaI")
                elif "news" in query:
                  from NewsRead import latestnews
                  latestnews()
                elif "send whatsapp" in query or "message on whatsapp" in query:
                 send_whatsapp_message_auto()
                
                elif "pause" in query:
                 pyautogui.press("k")
                 speak("video paused")
                elif "play" in query:
                 pyautogui.press("k")
                 speak("video played")
                elif "mute" in query:
                 pyautogui.press("m")
                 speak("video muted")
                elif "shutdown" in query or "restart" in query or "log off" in query:
                  handle_system_command(query)
                elif "youtube automation" in query:
                    speak("What YouTube action should I perform?")
                    action = takeCommand()
                    control_youtube(action)
                elif "seek forward" in query:
                 control_youtube("seek forward")
                elif "seek backward" in query:
                    control_youtube("seek backward")
                elif "increase speed" in query:
                    control_youtube("increase speed")
                elif "decrease speed" in query:
                    control_youtube("decrease speed")

                elif "gesture control" in query or "enable gesture control" in query:
                        speak("Activating gesture control mode. Here's what you can do:")
                        print("""\n
                        ================= GESTURE CONTROL MENU =================
                        Finger Count |     Action
                        -------------|-----------------------
                            1        | Play / Pause
                            2        | Volume Up
                            3        | Volume Down
                            4        | Scroll Up
                            5        | Scroll Down
                        ========================================================
                        """)
                        start_gesture_control()
                        speak("Gesture control session ended.")

                elif "volume up" in query:
                 from my_keyboard import volumeup
                 speak("Turning volume up,sir")
                 volumeup()
                elif "volume down" in query:
                 from my_keyboard import volumedown
                 speak("Turning volume down, sir")
                 volumedown()
                elif "scroll up" in query:
                  speak("scrolling up sir")
                  scroll_up(20)

                elif "scroll down" in query:
                  speak("scrolling down sir")
                  scroll_down(20)

                elif "open" in query:
                  from Dictapp import openappweb
                  openappweb(query)
                elif "close" in query:
                   from Dictapp import closeappweb
                   closeappweb(query)

                elif "google" in query:
                 from SearchNow import searchGoogle
                 searchGoogle(query)
                elif "youtube" in query:
                 from SearchNow import searchYoutube
                 searchYoutube(query)
                elif "wikipedia" in query:
                  from SearchNow import searchWikipedia
                  searchWikipedia(query)
                # elif "temperature" in query:
                #  search = "temperature in delhi"
                #  url = f"https://www.google.com/search?q={search}"
                #  r  = requests.get(url)
                #  data = BeautifulSoup(r.text,"html.parser")
                #  temp = data.find("div", class_ = "BNeawe").text
                #  speak(f"current{search} is {temp}")
                
                # elif "weather" in query:
                #   search = "temperature in delhi"
                #   url = f"https://www.google.com/search?q={search}"
                #   r  = requests.get(url)
                #   data = BeautifulSoup(r.text,"html.parser")
                #   temp = data.find("div", class_ = "BNeawe").text
                #   speak(f"current{search} is {temp}") 
                elif "temperature" in query or "weather" in query:
                    speak("Which city's weather do you want to know?")
                    city_query = takeCommand()
                    if city_query != "none":
                        get_weather(city_query)  # Call function with user-specified city
                    else:
                        get_weather()  # Default to Delhi if no input
                elif "battery" in query or "charge" in query or "battery percentage" in query:
                     check_battery()
                elif "set an alarm" in query:
                  print("input time example:- 10 and 10 and 10")
                  speak("Set the time")
                  a = input("Please tell the time :- ")
                  alarm(a)
                  speak("Done,sir")               
                elif "the time" in query:
                    strTime = datetime.datetime.now().strftime("%H:%M")    
                    speak(f"Sir, the time is {strTime}")
                elif "finally sleep" in query:
                 speak("Going to sleep,sir")
                 exit()
                elif "remember that" in query:
                 rememberMessage = query.replace("remember that","")
                 rememberMessage = query.replace("jarvis","")
                 speak("You told me to remember that"+rememberMessage)
                 remember = open("Remember.txt","a")
                 remember.write(rememberMessage)
                 remember.close()
                elif "what do you remember" in query:
                 remember = open("Remember.txt","r")
                 speak("You told me to remember that" + remember.read())
                elif "jarvis bot" in query or "chat bot" in query:
                  speak("Hello! I am your Jarvis ChatBot. Ask me anything.")
                  while True:
                     user_query = takeCommand()
                     if "exit chatbot" in user_query or "bye" in user_query:
                       speak("Exiting ChatBot mode. Back to main commands.")
                       break
                     elif user_query == "none":
                      continue
                    #  else:
                    #    response = chat_with_gpt(user_query)
                    #    print("Jarvis ChatBot:", response)
                    #    speak(response)
                     elif "make a pdf" in user_query or "create a pdf" in user_query:
                         if last_gpt_response:
                           create_pdf(last_gpt_response)
                         else:
                          speak("I don't have anything to save yet. Ask something first.")

                     elif "make a word" in user_query or "create word" in user_query:
                          if last_gpt_response:
                             create_word_doc(last_gpt_response)
                          else:
                              speak("I don't have anything to save yet. Ask something first.")

                     elif "make a ppt" in user_query or "create presentation" in user_query:
                          if last_gpt_response:
                           create_ppt(last_gpt_response)
                          else:
                           speak("I don't have anything to save yet. Ask something first.")
                     elif "generate content" in query:
                            speak("What content do you want me to generate?")
                            content_prompt = takeCommand()
                            result = chat_with_gpt(content_prompt)
                            print("Jarvis:", result)
                            speak(result)

                     else:
                        response = chat_with_gpt(user_query)
                        last_gpt_response = response  # Save for export
                        print("Jarvis ChatBot:", response)
                        speak(response)
                elif "generate image" in query or "make an image" in query:
                    speak("What image do you want me to generate?")
                    image_prompt = takeCommand()
                    generate_image(image_prompt)


                elif "calculate" in query:
                   from Calculatenumbers import WolfRamAlpha
                   from Calculatenumbers import Calc
                   query = query.replace("calculate","")
                   query = query.replace("jarvis","")
                   Calc(query)
                elif "screenshot" in query:  # <--- PASTE IT HERE
                   speak("Taking a screenshot now.")
                   time.sleep(1)
                   screenshot = pyautogui.screenshot()
                   screenshot.save("screenshot.jpg")
                   speak("Screenshot saved successfully.")
                   print("Screenshot saved as screenshot.jpg")
                elif "send email to" in query:
                    print("\n=== EMAIL COMPOSER ===")
                    recipient = input("Recipient Email: ").strip()
                    if not recipient:
                        speak("Email cancelled.")
                        continue
                    subject = input("Subject: ").strip()
                    if not subject:
                        speak("Email cancelled.")
                        continue
                    print("Body (press Enter twice to finish):")
                    body_lines = []
                    while True:
                        line = input()
                        if line == "":
                            break
                        body_lines.append(line)
                    body = "\n".join(body_lines)
                    if not body:
                        speak("Email cancelled.")
                        continue
                    speak("Sending email...")
                    success = email_assistant.send_email(recipient, subject, body)
                    if success:
                        speak(f"Email sent to {recipient}.")
                    else:
                        speak("Failed to send email. Please check your setup.")
                elif "read my emails" in query:
                    print("DEBUG: Command recognized - read my emails")
                    speak("Fetching your latest emails.")
                    emails = email_assistant.read_emails()
                    print(f"DEBUG: Emails fetched: {emails}")
                    if not emails:
                        print("DEBUG: No emails found or failed to fetch emails.")
                        speak("No emails found or failed to fetch emails.")
                    else:
                        for mail in emails:
                            print(f"DEBUG: Speaking email: From: {mail['from']}, Subject: {mail['subject']}, Body: {mail['body'][:100]}")
                            speak(f"From: {mail['from']}. Subject: {mail['subject']}. {mail['body'][:100]}")
                elif "add calendar event" in query:
                    print("\n=== CALENDAR EVENT COMPOSER ===")
                    title = input("Event Title: ").strip()
                    if not title:
                        speak("Event cancelled.")
                        continue
                    date = input("Date (MM/DD/YYYY): ").strip()
                    if not date:
                        speak("Event cancelled.")
                        continue
                    time = input("Time (HH:MM) or press Enter for all day: ").strip()
                    description = input("Description (optional): ").strip()
                    speak("Adding calendar event...")
                    success = calendar_assistant.add_event(title, date, time, description)
                    if success:
                        speak(f"Event '{title}' added to calendar.")
                    else:
                        speak("Failed to add event. Please check the date format.")
                elif "list calendar events" in query:
                    speak("Fetching your calendar events.")
                    events = calendar_assistant.list_events()
                    if events:
                        speak(f"Found {len(events)} upcoming events.")
                        for event in events:
                            event_date = datetime.fromisoformat(event['datetime'])
                            speak(f"Event: {event['title']} on {event_date.strftime('%B %d at %I:%M %p')}")
                    else:
                        speak("No upcoming events found.")
                elif "today's events" in query or "today events" in query:
                    speak("Checking today's events.")
                    today_events = calendar_assistant.get_today_events()
                    if today_events:
                        speak(f"You have {len(today_events)} events today.")
                        for event in today_events:
                            event_date = datetime.fromisoformat(event['datetime'])
                            speak(f"Event: {event['title']} at {event_date.strftime('%I:%M %p')}")
                    else:
                        speak("No events scheduled for today.")
                elif "open calendar" in query or "open system calendar" in query:
                    speak("Opening system calendar.")
                    calendar_assistant.open_system_calendar()
                elif "open google calendar" in query or "open gmail calendar" in query:
                    speak("Opening Google Calendar.")
                    calendar_assistant.open_google_calendar()
                elif "add event to google calendar" in query or "add to google calendar" in query:
                    print("\n=== GOOGLE CALENDAR EVENT ===")
                    title = input("Event Title: ").strip()
                    if not title:
                        speak("Event cancelled.")
                        continue
                    date = input("Date (MM/DD/YYYY): ").strip()
                    if not date:
                        speak("Event cancelled.")
                        continue
                    time = input("Time (HH:MM) or press Enter for all day: ").strip()
                    description = input("Description (optional): ").strip()
                    speak("Opening Google Calendar with your event.")
                    success = calendar_assistant.open_calendar_with_event(title, date, time, description)
                    if success:
                        speak(f"Google Calendar opened with event '{title}'.")
                    else:
                        speak("Failed to open Google Calendar.")
                elif "add reminder" in query:
                    print("\n=== REMINDER COMPOSER ===")
                    title = input("Reminder Title: ").strip()
                    if not title:
                        speak("Reminder cancelled.")
                        continue
                    due_date = input("Due Date (MM/DD/YYYY) or press Enter for no due date: ").strip()
                    priority = input("Priority (high/medium/low) [default: medium]: ").strip() or "medium"
                    description = input("Description (optional): ").strip()
                    speak("Adding reminder...")
                    success = reminders_assistant.add_reminder(title, due_date, priority, description)
                    if success:
                        speak(f"Reminder '{title}' added successfully.")
                    else:
                        speak("Failed to add reminder. Please check the date format.")
                elif "add todo" in query:
                    print("\n=== TODO COMPOSER ===")
                    title = input("Todo Title: ").strip()
                    if not title:
                        speak("Todo cancelled.")
                        continue
                    priority = input("Priority (high/medium/low) [default: medium]: ").strip() or "medium"
                    description = input("Description (optional): ").strip()
                    speak("Adding todo item...")
                    success = reminders_assistant.add_todo(title, priority, description)
                    if success:
                        speak(f"Todo '{title}' added successfully.")
                    else:
                        speak("Failed to add todo item.")
                elif "list reminders" in query:
                    speak("Fetching your reminders.")
                    reminders = reminders_assistant.list_reminders()
                    if reminders:
                        speak(f"Found {len(reminders)} active reminders.")
                        for reminder in reminders:
                            speak(f"Reminder: {reminder['title']} - Priority: {reminder['priority']}")
                    else:
                        speak("No active reminders found.")
                elif "list todos" in query:
                    speak("Fetching your todo list.")
                    todos = reminders_assistant.list_todos()
                    if todos:
                        speak(f"Found {len(todos)} active todo items.")
                        for todo in todos:
                            speak(f"Todo: {todo['title']} - Priority: {todo['priority']}")
                    else:
                        speak("No active todo items found.")
                elif "overdue reminders" in query or "check overdue" in query:
                    speak("Checking for overdue reminders.")
                    overdue = reminders_assistant.get_overdue_reminders()
                    if overdue:
                        speak(f"Found {len(overdue)} overdue reminders.")
                        for reminder in overdue:
                            speak(f"Overdue: {reminder['title']}")
                    else:
                        speak("No overdue reminders found.")
                elif "test notification" in query:
                    speak("Testing notification system.")
                    notification_system.test_notification_system()
                    speak("Test notification sent.")
                elif "start notifications" in query:
                    speak("Starting notification service.")
                    reminders_assistant.start_notification_service()
                    speak("Notification service is now active.")
                elif "stop notifications" in query:
                    speak("Stopping notification service.")
                    reminders_assistant.stop_notification_service()
                    speak("Notification service stopped.")
                
                # Advanced Features
                elif "system monitor" in query or "start monitoring" in query:
                    speak("Starting system monitoring.")
                    result = advanced_features.system_monitor()
                    speak(result)
                
                elif "system status" in query or "check system" in query:
                    speak("Checking system status.")
                    status = advanced_features.get_system_status()
                    print(status)
                    speak("System status displayed on screen.")
                
                elif "encrypt file" in query:
                    speak("Please specify the file path to encrypt.")
                    file_path = input("File path: ").strip()
                    if file_path and os.path.exists(file_path):
                        password = input("Enter password: ").strip()
                        if password:
                            result = advanced_features.encrypt_file(file_path, password)
                            speak(result)
                        else:
                            speak("Encryption cancelled.")
                    else:
                        speak("File not found or path invalid.")
                
                elif "decrypt file" in query:
                    speak("Please specify the file path to decrypt.")
                    file_path = input("File path: ").strip()
                    if file_path and os.path.exists(file_path):
                        password = input("Enter password: ").strip()
                        if password:
                            result = advanced_features.decrypt_file(file_path, password)
                            speak(result)
                        else:
                            speak("Decryption cancelled.")
                    else:
                        speak("File not found or path invalid.")
                
                elif "generate qr code" in query:
                    speak("What data should I encode in the QR code?")
                    data = takeCommand()
                    if data != "none":
                        result = advanced_features.generate_qr_code(data)
                        speak(result)
                    else:
                        speak("QR code generation cancelled.")
                
                elif "read qr code" in query:
                    speak("Please specify the image path.")
                    image_path = input("Image path: ").strip()
                    if image_path and os.path.exists(image_path):
                        result = advanced_features.read_qr_code(image_path)
                        speak(f"QR code contains: {result}")
                    else:
                        speak("Image not found or path invalid.")
                
                elif "generate password" in query:
                    speak("Generating secure password.")
                    password = advanced_features.password_generator()
                    speak(f"Generated password: {password}")
                    print(f"🔐 Generated Password: {password}")
                
                elif "organize files" in query:
                    speak("Please specify the directory to organize.")
                    directory = input("Directory path: ").strip()
                    if directory and os.path.exists(directory):
                        result = advanced_features.file_organizer(directory)
                        speak(result)
                    else:
                        speak("Directory not found or path invalid.")
                
                elif "analysis" in query:
                    speak("What text should I analyze?")
                    text = takeCommand()
                    if text != "none":
                        speak("Analyzing text with AI.")
                        result = advanced_features.ai_analysis(text)
                        print("🤖 AI Analysis:", result)
                        speak("Analysis complete. Check the screen for details.")
                    else:
                        speak("Analysis cancelled.")
                
                elif "create backup" in query:
                    speak("Please specify source and backup paths.")
                    source = input("Source path: ").strip()
                    backup = input("Backup path: ").strip()
                    if source and backup and os.path.exists(source):
                        result = advanced_features.create_backup(source, backup)
                        speak(result)
                    else:
                        speak("Invalid paths or source not found.")
                
                elif "scan network" in query:
                    speak("Scanning network for devices.")
                    result = advanced_features.network_scanner()
                    if isinstance(result, list):
                        speak(f"Found {len(result)} devices on network.")
                        for device in result:
                            print(f"📱 Device: {device['hostname']} ({device['ip']})")
                    else:
                        speak(result)
                
                elif "test advanced notifications" in query:
                    speak("Testing advanced notification system.")
                    notification_system.test_notification_system()
                    speak("Advanced notification test complete.")
                
                elif "configure notifications" in query:
                    speak("Configuring notification settings.")
                    sound = input("Enable sound? (y/n): ").lower() == 'y'
                    style = input("Style (default/modern/minimal): ").strip() or "default"
                    auto_open = input("Auto-open links? (y/n): ").lower() == 'y'
                    notification_system.configure_notifications(sound, style, auto_open)
                    speak("Notification settings updated.")
                
                elif "show notification history" in query:
                    speak("Showing notification history.")
                    history = notification_system.get_notification_history()
                    if history:
                        print("📋 Notification History:")
                        for record in history[-10:]:  # Show last 10
                            print(f"  {record['timestamp']}: {record['title']} - {record['message']}")
                    else:
                        print("No notification history found.")
                    speak("Notification history displayed.")
                
                elif "clear notification history" in query:
                    speak("Clearing notification history.")
                    notification_system.clear_notification_history()
                    speak("Notification history cleared.")
                
                # AI Personality Features
                elif "switch personality" in query:
                    speak("Which personality would you like to switch to?")
                    personality = takeCommand()
                    if personality != "none":
                        result = personality_manager.switch_personality(personality)
                        speak(result)
                    else:
                        speak("Personality switch cancelled.")
                
                elif "personality info" in query:
                    info = personality_manager.get_personality_info()
                    speak(info)
                
                elif "list personalities" in query:
                    personalities = personality_manager.list_personalities()
                    speak("Available personalities:")
                    for personality in personalities:
                        print(f"  • {personality}")
                
                # Analytics Dashboard Features
                elif "show my stats" in query:
                    speak("Generating comprehensive analytics report.")
                    report = analytics_dashboard.generate_detailed_report()
                    print(report)
                    speak("Analytics report displayed on screen.")
                
                elif "usage statistics" in query:
                    speak("Fetching usage statistics.")
                    stats = analytics_dashboard.get_usage_stats()
                    if stats:
                        print("📊 Usage Statistics (Last 7 Days):")
                        print(f"  • Total Commands: {stats['total_commands']}")
                        print(f"  • Successful Commands: {stats['successful_commands']}")
                        success_rate = (stats['successful_commands']/stats['total_commands']*100) if stats['total_commands'] > 0 else 0
                        print(f"  • Success Rate: {success_rate:.1f}%")
                        print(f"  • Most Used Command: {stats['most_used_command']}")
                        print(f"  • Average Response Time: {stats['avg_response_time']:.2f}s")
                    else:
                        print("No usage statistics available.")
                    speak("Usage statistics displayed.")
                
                elif "system performance report" in query:
                    speak("Generating system performance report.")
                    report = analytics_dashboard.generate_performance_report()
                    print(report)
                    speak("System performance report displayed.")
                
                # Voice Games Features
                elif "start trivia" in query:
                    speak("Starting trivia game. Get ready for some questions!")
                    result = voice_games.start_trivia()
                    speak(result)
                
                elif "next question" in query:
                    result = voice_games.get_next_question()
                    if result:
                        speak(result)
                    else:
                        speak("No trivia game active. Say 'start trivia' to begin.")
                
                elif "stop trivia" in query:
                    result = voice_games.stop_trivia()
                    speak(result)
                
                elif "start word game" in query:
                    speak("Starting word guessing game. I'll give you clues!")
                    result = voice_games.start_word_game()
                    speak(result)
                
                elif "next clue" in query:
                    result = voice_games.get_next_clue()
                    if result:
                        speak(result)
                    else:
                        speak("No word game active. Say 'start word game' to begin.")
                
                elif "stop word game" in query:
                    result = voice_games.stop_word_game()
                    speak(result)
                
                elif "start story" in query:
                    speak("Starting interactive story generator. I'll ask you questions to create your story!")
                    result = voice_games.start_story()
                    speak(result)
                
                elif "tell story" in query:
                    result = voice_games.tell_story()
                    if result:
                        speak("Here's your story:")
                        print(f"📖 {result}")
                        speak("Story displayed on screen.")
                    else:
                        speak("No story available. Say 'start story' to begin.")
                
                elif "stop story" in query:
                    result = voice_games.stop_story()
                    speak(result)
                
                elif "game status" in query:
                    status = voice_games.get_game_status()
                    speak(status)
                
                elif "my answer is" in query or "the answer is" in query:
                    # Extract answer from query
                    if "my answer is" in query:
                        answer = query.replace("my answer is", "").strip()
                    else:
                        answer = query.replace("the answer is", "").strip()
                    
                    if answer:
                        result = voice_games.submit_answer(answer)
                        speak(result)
                    else:
                        speak("Please provide an answer.")
                
                elif "story answer" in query:
                    # Extract story element from query
                    answer = query.replace("story answer", "").strip()
                    if answer:
                        result = voice_games.add_story_element(answer)
                        speak(result)
                    else:
                        speak("Please provide a story element.")
                
                # Enhanced AI Features
                elif "ai chat" in query or "advanced chat" in query:
                    speak("What would you like to discuss with AI?")
                    message = takeCommand()
                    if message != "none":
                        speak("Processing with advanced AI.")
                        result = enhanced_ai.advanced_chat(message)
                        print("🤖 AI Response:", result)
                        speak("AI response displayed on screen.")
                    else:
                        speak("AI chat cancelled.")
                
                elif "generate ai image" in query or "create ai image" in query:
                    speak("What image should I generate?")
                    prompt = takeCommand()
                    if prompt != "none":
                        speak("Generating AI image.")
                        result = enhanced_ai.generate_image(prompt)
                        speak(result)
                    else:
                        speak("Image generation cancelled.")
                
                elif "analyze sentiment" in query:
                    speak("What text should I analyze for sentiment?")
                    text = takeCommand()
                    if text != "none":
                        speak("Analyzing sentiment.")
                        result = enhanced_ai.analyze_sentiment(text)
                        print("😊 Sentiment Analysis:", result)
                        speak("Sentiment analysis complete.")
                    else:
                        speak("Sentiment analysis cancelled.")
                
                elif "translate" in query:
                    speak("What language should I translate to?")
                    language = takeCommand()
                    if language != "none":
                        speak("What text should I translate?")
                        text = takeCommand()
                        if text != "none":
                            speak("Translating text.")
                            result = enhanced_ai.translate_text(text, language)
                            print("🌐 Translation:", result)
                            speak("Translation complete.")
                        else:
                            speak("Translation cancelled.")
                    else:
                        speak("Translation cancelled.")
                
                elif "summarize" in query:
                    speak("What text should I summarize?")
                    text = takeCommand()
                    if text != "none":
                        speak("Summarizing text.")
                        result = enhanced_ai.summarize_text(text)
                        print("📝 Summary:", result)
                        speak("Summarization complete.")
                    else:
                        speak("Summarization cancelled.")
                
                elif "code help" in query or "programming help" in query:
                    speak("What coding question do you have?")
                    question = takeCommand()
                    if question != "none":
                        speak("Analyzing your code question.")
                        result = enhanced_ai.code_assistant(question)
                        print("💻 Code Assistant:", result)
                        speak("Code assistance provided.")
                    else:
                        speak("Code help cancelled.")
                
                elif "smart reminder" in query:
                    speak("What task should I create a smart reminder for?")
                    task = takeCommand()
                    if task != "none":
                        speak("Creating smart reminder.")
                        result = enhanced_ai.smart_reminder(task)
                        print("📅 Smart Reminder:", result)
                        speak("Smart reminder created.")
                    else:
                        speak("Smart reminder cancelled.")
                
                elif "ai research" in query or "research topic" in query:
                    speak("What topic should I research?")
                    topic = takeCommand()
                    if topic != "none":
                        speak("Conducting AI research.")
                        result = enhanced_ai.ai_research(topic)
                        print("🔬 AI Research:", result)
                        speak("Research complete.")
                    else:
                        speak("Research cancelled.")
                
                elif "switch ai context" in query or "change ai mode" in query:
                    speak("Which AI context would you like? Available: assistant, creative, technical, educational, business")
                    context = takeCommand()
                    if context != "none":
                        result = enhanced_ai.set_context(context)
                        speak(result)
                    else:
                        speak("Context switch cancelled.")
                
                elif "ai status" in query or "ai info" in query:
                    status = enhanced_ai.get_ai_status()
                    speak(f"Current AI context: {status['current_context']}. Conversation length: {status['conversation_length']} messages.")
                    print("🤖 AI Status:", status)
                
                elif "analyze conversation" in query:
                    speak("Analyzing our conversation.")
                    result = enhanced_ai.conversation_analysis()
                    print("📊 Conversation Analysis:", result)
                    speak("Conversation analysis complete.")
                
                elif "clear ai history" in query:
                    result = enhanced_ai.clear_history()
                    speak(result)
                
                elif "text to speech" in query or "convert to speech" in query:
                    speak("What text should I convert to speech?")
                    text = takeCommand()
                    if text != "none":
                        speak("Converting text to speech.")
                        result = enhanced_ai.text_to_speech(text)
                        speak(result)
                    else:
                        speak("Text-to-speech cancelled.")
                
                # Track command for analytics
                else:
                    # Track unsuccessful command
                    analytics_dashboard.track_command(query, False)
                    speak("I didn't understand that command. Please try again.")
                
                # Track successful command
                analytics_dashboard.track_command(query, True)